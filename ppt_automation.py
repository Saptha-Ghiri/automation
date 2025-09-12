from pptx import Presentation
from pptx.chart.data import CategoryChartData
import re

def generate_weekly_report(
    template_path,
    output_path,
    report_date,
    new_period,
    total_tasks,
    completed_tasks,
    ticket_status_data,
    individual_data,
    main_chart_data,
    pie1_data,
    pie2_data,
    new_date,
    slide5_data=None,
    slide6_data=None
):
    prs = Presentation(template_path)

    # --- SLIDE 1: Update report date ---
    slide1 = prs.slides[0]
    date_pattern = re.compile(r"\d{1,2}\s+[A-Za-z]+\s+\d{4}")

    for shape in slide1.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                full_text = "".join(run.text for run in paragraph.runs)
                if date_pattern.search(full_text):
                    if paragraph.runs:
                        first_run = paragraph.runs[0]
                        for run in paragraph.runs:
                            run.text = ""
                        first_run.text = report_date

    # --- SLIDE 2: Update period, tasks, SLA ---
    prs = Presentation(template_path)
    slide2 = prs.slides[1]

    # Regex to find text inside parentheses
    pattern = re.compile(r"\(.*?\)")

    for shape in slide2.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                full_text = "".join(run.text for run in paragraph.runs)
                if "(" in full_text and ")" in full_text:
                    new_text = pattern.sub(f"({new_period})", full_text)
                    
                    # Clear old runs and insert new formatted text
                    if paragraph.runs:
                        first_run = paragraph.runs[0]
                        for run in paragraph.runs:
                            run.text = ""
                        first_run.text = new_text
    

    # Calculate percentages
    percent_completed = (completed_tasks / total_tasks * 100) if total_tasks > 0 else 0
    sla_percent = (completed_tasks / total_tasks * 100) if total_tasks > 0 else 0

    for shape in slide2.shapes:
        if shape.has_table:
            table = shape.table

            # Assume first row = headers, second row = values
            headers = [cell.text.strip() for cell in table.rows[0].cells]

            for i, header in enumerate(headers):
                if header.startswith("Total Task"):
                    table.cell(1, i).text = str(total_tasks)

                elif header.startswith("Task Completed"):
                    table.cell(1, i).text = str(completed_tasks)

                elif header.startswith("% Completed"):
                    table.cell(1, i).text = f"{percent_completed:.0f}%"

                elif header.startswith("Met SLA %"):
                    # Update header with (x/y)
                    table.cell(0, i).text = f"Met SLA % ({completed_tasks}/{total_tasks})"
                    table.cell(1, i).text = f"{sla_percent:.0f}%"
                        
    # --- SLIDE 2: Update charts ---
    charts = [shape.chart for shape in slide2.shapes if shape.has_chart]

    if len(charts) >= 2:
        chart1, chart2 = charts[0], charts[1]

        # Chart 1: Ticket Status
        chart_data1 = CategoryChartData()
        chart_data1.categories = list(ticket_status_data.keys())
        chart_data1.add_series("Ticket Status", list(ticket_status_data.values()))
        chart1.replace_data(chart_data1)

        # Chart 2: Individual Completion
        chart_data2 = CategoryChartData()
        chart_data2.categories = list(individual_data.keys())
        chart_data2.add_series("Completed Tasks", list(individual_data.values()))
        chart2.replace_data(chart_data2)

    # --- SLIDE 3: Update charts ---
    slide3 = prs.slides[2]  # Slide 3 (index 2)

    charts = [shape.chart for shape in slide3.shapes if shape.has_chart]

    if len(charts) >= 3:
        chart_main, chart_pie1, chart_pie2 = charts[0], charts[1], charts[2]

        # --- 1. Update main chart ---
        chart_data_main = CategoryChartData()
        chart_data_main.categories = list(main_chart_data.keys())
        # Keep existing series name
        series_name_main = chart_main.series[0].name if chart_main.series else "Main Chart"
        chart_data_main.add_series(series_name_main, list(main_chart_data.values()))
        chart_main.replace_data(chart_data_main)

        # --- 2. Update first pie chart ---
        pie_data1 = CategoryChartData()
        pie_data1.categories = list(pie1_data.keys())
        series_name_pie1 = chart_pie1.series[0].name if chart_pie1.series else "Pie 1"
        pie_data1.add_series(series_name_pie1, list(pie1_data.values()))
        chart_pie1.replace_data(pie_data1)

        # --- 3. Update second pie chart ---
        pie_data2 = CategoryChartData()
        pie_data2.categories = list(pie2_data.keys())
        series_name_pie2 = chart_pie2.series[0].name if chart_pie2.series else "Pie 2"
        pie_data2.add_series(series_name_pie2, list(pie2_data.values()))
        chart_pie2.replace_data(pie_data2)

    #---slide 4: Update Date---
    slide4 = prs.slides[3]  # Slide 4 (index 3)

    # Loop through all shapes on the slide

    for shape in slide4.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                # Check if line contains 'Date:'
                if 'Date:' in paragraph.text:
                    # Replace only the date after 'Date:'
                    paragraph.text = re.sub(r'Date:\s*\d{2}/\d{2}/\d{4}', f'Date: {new_date}', paragraph.text)

    #---SLIDE 5: DaaS Queue Monitoring Update---
    if len(prs.slides) >= 5:
        slide5 = prs.slides[4]  # Slide 5 (index 4)

        # Update text content in slide 5 based on slide5_data
        if slide5_data:
            daily_shapes = []
            for shape in slide5.shapes:
                if shape.has_text_frame:
                    # Shape 1: Update period in header
                    if "DaaS - Queue Monitoring status for the period of" in shape.text_frame.text:
                        for paragraph in shape.text_frame.paragraphs:
                            full_text = "".join(run.text for run in paragraph.runs)
                            if "(" in full_text and ")" in full_text:
                                pattern = re.compile(r"\(.*?\)")
                                new_text = pattern.sub(f"({new_period})", full_text)
                                if paragraph.runs:
                                    first_run = paragraph.runs[0]
                                    for run in paragraph.runs:
                                        run.text = ""
                                    first_run.text = new_text
                    
                    # Shape 2: Update summary statistics
                    elif "Total No. of Tickets" in shape.text_frame.text:
                        summary_stats = slide5_data.get('summary_stats', {})
                        paragraphs = list(shape.text_frame.paragraphs)
                        
                        # Update existing paragraphs to preserve formatting
                        if len(paragraphs) >= 1 and "Total No. of Tickets" in paragraphs[0].text:
                            new_text = f"Total No. of Tickets : {summary_stats.get('total_tickets', 155)}"
                            if paragraphs[0].runs:
                                paragraphs[0].runs[0].text = new_text
                                for i in range(1, len(paragraphs[0].runs)):
                                    paragraphs[0].runs[i].text = ""
                            else:
                                paragraphs[0].text = new_text
                                
                        if len(paragraphs) >= 2 and "Awaiting" in paragraphs[1].text:
                            new_text = f"Awaiting : {summary_stats.get('awaiting', 5):02d}"
                            if paragraphs[1].runs:
                                paragraphs[1].runs[0].text = new_text
                                for i in range(1, len(paragraphs[1].runs)):
                                    paragraphs[1].runs[i].text = ""
                            else:
                                paragraphs[1].text = new_text
                                
                        if len(paragraphs) >= 3 and "Ticket Closed" in paragraphs[2].text:
                            new_text = f"Ticket Closed : {summary_stats.get('closed', 2):02d}"
                            if paragraphs[2].runs:
                                paragraphs[2].runs[0].text = new_text
                                for i in range(1, len(paragraphs[2].runs)):
                                    paragraphs[2].runs[i].text = ""
                            else:
                                paragraphs[2].text = new_text
                                
                        if len(paragraphs) >= 4 and "Resolved with Customer" in paragraphs[3].text:
                            new_text = f"Resolved with Customer : {summary_stats.get('resolved', 148)}"
                            if paragraphs[3].runs:
                                paragraphs[3].runs[0].text = new_text
                                for i in range(1, len(paragraphs[3].runs)):
                                    paragraphs[3].runs[i].text = ""
                            else:
                                paragraphs[3].text = new_text
                    
                    # Collect shapes with daily data
                    elif any(f"Date: {date}" in shape.text_frame.text for date in ["09/01/2025", "09/02/2025", "09/03/2025", "09/04/2025", "0905/2025"]):
                        daily_shapes.append(shape)
            
            # Update daily shapes with proper date distribution
            daily_data = slide5_data.get('daily_data', {})
            if daily_data and daily_shapes:
                daily_dates = list(daily_data.keys())
                
                # Shape 4 gets first 3 dates, Shape 5 gets last 2 dates
                for shape_idx, shape in enumerate(daily_shapes):
                    paragraphs = list(shape.text_frame.paragraphs)
                    
                    if shape_idx == 0:  # Shape 4 - first 3 dates
                        dates_for_shape = daily_dates[:3]
                    elif shape_idx == 1:  # Shape 5 - last 2 dates  
                        dates_for_shape = daily_dates[3:5]
                    else:
                        continue
                    
                    current_date_index = 0
                    current_date_key = dates_for_shape[0] if dates_for_shape else None
                    
                    # Update existing paragraphs to preserve formatting
                    for i, paragraph in enumerate(paragraphs):
                        if "Date:" in paragraph.text:
                            # Use the next date for this shape
                            if current_date_index < len(dates_for_shape):
                                current_date_key = dates_for_shape[current_date_index]
                                new_text = f"Date: {current_date_key}:"
                                if paragraph.runs:
                                    paragraph.runs[0].text = new_text
                                    for j in range(1, len(paragraph.runs)):
                                        paragraph.runs[j].text = ""
                                else:
                                    paragraph.text = new_text
                                current_date_index += 1
                        elif "No of Tickets by" in paragraph.text and current_date_key:
                            # Extract person name and update ticket count for current date
                            person_match = re.search(r'No of Tickets by (\w+)', paragraph.text)
                            if person_match:
                                person = person_match.group(1)
                                people_data = daily_data.get(current_date_key, {})
                                if person in people_data:
                                    new_text = f"No of Tickets by {person} - {people_data[person]:02d}"
                                    if paragraph.runs:
                                        paragraph.runs[0].text = new_text
                                        for j in range(1, len(paragraph.runs)):
                                            paragraph.runs[j].text = ""
                                    else:
                                        paragraph.text = new_text
        else:
            # Fallback: Just update dates and periods like other slides
            for shape in slide5.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        # Update date references
                        if 'Date:' in paragraph.text:
                            paragraph.text = re.sub(r'Date:\s*\d{2}/\d{2}/\d{4}', f'Date: {new_date}', paragraph.text)
                        
                        # Update period references
                        full_text = "".join(run.text for run in paragraph.runs)
                        if "(" in full_text and ")" in full_text:
                            pattern = re.compile(r"\(.*?\)")
                            new_text = pattern.sub(f"({new_period})", full_text)
                            
                            if paragraph.runs:
                                first_run = paragraph.runs[0]
                                for run in paragraph.runs:
                                    run.text = ""
                                first_run.text = new_text

    #---SLIDE 6: Update charts (if slide 6 exists)---
    if len(prs.slides) >= 6:
        slide6 = prs.slides[5]  # Slide 6 (index 5)

        # Update text content in slide 6 (dates and periods)
        for shape in slide6.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    # Update date references
                    if 'Date:' in paragraph.text:
                        new_text = re.sub(r'Date:\s*\d{2}/\d{2}/\d{4}', f'Date: {new_date}', paragraph.text)
                        if paragraph.runs:
                            paragraph.runs[0].text = new_text
                            for j in range(1, len(paragraph.runs)):
                                paragraph.runs[j].text = ""
                        else:
                            paragraph.text = new_text
                    
                    # Update period references
                    full_text = "".join(run.text for run in paragraph.runs)
                    if "(" in full_text and ")" in full_text:
                        pattern = re.compile(r"\(.*?\)")
                        new_text = pattern.sub(f"({new_period})", full_text)
                        
                        if paragraph.runs:
                            paragraph.runs[0].text = new_text
                            for j in range(1, len(paragraph.runs)):
                                paragraph.runs[j].text = ""
                        else:
                            paragraph.text = new_text

        # Update slide 6 content based on template structure
        if slide6_data:
            for shape in slide6.shapes:
                if shape.has_text_frame:
                    # Update header text with period
                    if "DaaS - Queue Monitoring comparison" in shape.text_frame.text:
                        for paragraph in shape.text_frame.paragraphs:
                            if "comparison with previous week" in paragraph.text:
                                new_text = f"DaaS - Queue Monitoring comparison with previous week ({new_period})"
                                if paragraph.runs:
                                    paragraph.runs[0].text = new_text
                                    for j in range(1, len(paragraph.runs)):
                                        paragraph.runs[j].text = ""
                                else:
                                    paragraph.text = new_text
                    
                    # Update any date references
                    for paragraph in shape.text_frame.paragraphs:
                        if 'Date:' in paragraph.text:
                            new_text = re.sub(r'Date:\s*\d{2}/\d{2}/\d{4}', f'Date: {new_date}', paragraph.text)
                            if paragraph.runs:
                                paragraph.runs[0].text = new_text
                                for j in range(1, len(paragraph.runs)):
                                    paragraph.runs[j].text = ""
                            else:
                                paragraph.text = new_text

            # Update charts with sample data (if charts allow it)
            charts = [shape.chart for shape in slide6.shapes if shape.has_chart]
            
            if len(charts) >= 2:
                # Chart 1: Column chart - Weekly comparison data
                if 'column_chart_data' in slide6_data:
                    try:
                        chart1 = charts[0]
                        chart_data1 = CategoryChartData()
                        chart_data1.categories = list(slide6_data['column_chart_data'].keys())
                        chart_data1.add_series("Weekly Tickets", list(slide6_data['column_chart_data'].values()))
                        chart1.replace_data(chart_data1)
                    except Exception as e:
                        pass  # Chart may have external data source
                
                # Chart 2: Bar chart - Status breakdown (3 series)
                if 'bar_chart_data' in slide6_data:
                    try:
                        chart2 = charts[1]
                        chart_data2 = CategoryChartData()
                        bar_data = slide6_data['bar_chart_data']
                        chart_data2.categories = list(bar_data.keys())
                        chart_data2.add_series("Awaiting", [data['awaiting'] for data in bar_data.values()])
                        chart_data2.add_series("Ticket Closed", [data['closed'] for data in bar_data.values()])
                        chart_data2.add_series("Resolved with Customer", [data['resolved'] for data in bar_data.values()])
                        chart2.replace_data(chart_data2)
                    except Exception as e:
                        pass  # Chart may have external data source

            # Update tables if present
            if 'tables' in slide6_data:
                table_shapes = [shape for shape in slide6.shapes if shape.has_table]
                for i, shape in enumerate(table_shapes):
                    if i < len(slide6_data['tables']):
                        table = shape.table
                        table_data = slide6_data['tables'][i]
                        # Update specific table cells
                        for row_idx, row_data in enumerate(table_data.get('rows', [])):
                            if row_idx < len(table.rows):
                                for col_idx, cell_value in enumerate(row_data):
                                    if col_idx < len(table.columns):
                                        table.cell(row_idx, col_idx).text = str(cell_value)

    # --- Save final presentation ---
    prs.save(output_path)


# ------------------------
# ✅ Example usage
# ------------------------
ticket_status_data = {
    "resolvedwith customer": 22,
    "internal solution provided": 11,
    "awaiting": 30,
    "inprogress": 70,
    "new": 20
}

individual_data = {
    "Abhijeet": 14,
    "Aditya": 61,
    "Nishanth": 57,
    "Sakthivel": 12
}

main_chart_data = {
    "Atomic" :	10,
    "Beigene"	: 18,
    "BMS"	: 10,
    "Collegum"	: 10,
    "Azure Imdaas" :	11,
    "AWS Imdaas"	: 12,
    "MDM"	: 10,
    "Usbu-Pede"	: 10,
}

pie1_data = {
    "SLA Met" :	100,
    "SLA Lost"	: 10

}

pie2_data = {
"Priority 1" :	10,
"Priority 2" :	10,
"Priority 3" :	19,
"Priority 4" :	92
}

new_date = "09/09/2025"

# Slide 6 data structure based on template analysis
slide6_data = {
    'column_chart_data': {
        "Augest 2nd Week": 180,
        "August 3rd Week": 200, 
        "August 4th Week": 220,
        "September 1st Week": 155
    },
    'bar_chart_data': {
        "Augest 2nd Week": {
            'awaiting': 8,
            'closed': 5,
            'resolved': 187
        },
        "August 3rd Week": {
            'awaiting': 12,
            'closed': 3,
            'resolved': 165
        },
        "August 4th Week": {
            'awaiting': 12,
            'closed': 3,
            'resolved': 165
        },
        "September 1st Week": {
            'awaiting': 12,
            'closed': 3,
            'resolved': 165
        }
    }
}

# Slide 5 data structure for DaaS Queue Monitoring
slide5_data = {
    'summary_stats': {
        'total_tickets': 200,
        'awaiting': 8,
        'closed': 5,
        'resolved': 187
    },
    'daily_data': {
        '09/01/2025': {
            'Abhjieet': 35,
            'Saptha': 5,
            'Sakthivel': 8,
            'Aditya': 2
        },
        '09/02/2025': {
            'Abhjieet': 25,
            'Saptha': 6,
            'Sakthivel': 12,
            'Aditya': 1
        },
        '09/03/2025': {
            'Abhjieet': 30,
            'Saptha': 4,
            'Sakthivel': 7,
            'Aditya': 3
        },
        '09/04/2025': {
            'Abhjieet': 28,
            'Saptha': 2,
            'Sakthivel': 9,
            'Aditya': 1
        },
        '09/05/2025': {
            'Abhjieet': 22,
            'Saptha': 4,
            'Sakthivel': 15,
            'Aditya': 2
        }
    }
}

generate_weekly_report(
    "template.pptx",
    "final_report.pptx",
    report_date="19 September 2025",
    new_period="09/01/2025 to 09/09/2025",
    total_tasks=27,
    completed_tasks=23,
    ticket_status_data=ticket_status_data,
    individual_data=individual_data,
    main_chart_data = main_chart_data,
    pie1_data = pie1_data,
    pie2_data = pie2_data, 
    new_date = new_date,
    slide5_data = slide5_data,
    slide6_data = slide6_data  # Add slide 6 parameter
)
